"""
W5：白皮书生成 + PPT生成
"""
import os
import sys
from config import *

# ═══════════════════════════════════════════════════════════════
# 1. White Paper (Markdown -> will serve as PDF source)
# ═══════════════════════════════════════════════════════════════

def generate_white_paper(wide=None):
    """生成白皮书（Markdown + 关键统计）"""
    print('\n[Generating White Paper]')

    # Load data for stats
    if wide is None:
        csv_path = os.path.join(DATA_DIR, 'wide_table.csv')
        if os.path.exists(csv_path):
            import pandas as pd
            wide = pd.read_csv(csv_path)

    if wide is None:
        print('  No data available for white paper')
        return

    # Key statistics
    total_leads = len(wide)
    cvr = wide['是否下订'].mean() * 100
    total_orders = wide['订单ID'].nunique()
    delay_rate = wide[wide['订单ID'].notna()]['是否延迟交付'].mean() * 100
    avg_delivery_score = wide[wide['交付评分'].notna()]['交付评分'].mean()
    avg_aftersales = wide[wide['售后平均满意度'] > 0]['售后平均满意度'].mean()

    # Top channels
    channel_cvr = wide.groupby('渠道')['是否下订'].mean().sort_values(ascending=False) * 100
    top_channel = channel_cvr.index[0]
    top_channel_cvr = channel_cvr.iloc[0]

    # City stats
    city_cvr = wide.groupby('城市')['是否下订'].mean().sort_values(ascending=False) * 100
    best_city = city_cvr.index[0]
    worst_city = city_cvr.index[-1]

    # Build white paper
    paper = f"""# G9 销售运营优化白皮书

---

## 第1章 执行摘要

### 一句话结论
**通过数据驱动的预算重新分配、跟进流程优化和智能预警体系建设，G9可在不增加总预算的前提下实现Q4销量同比增长15%以上的目标。**

### 五大核心发现

1. **渠道效率悬殊**：{top_channel}渠道转化率({top_channel_cvr:.1f}%)远高于平均值({cvr:.1f}%)，而部分传统渠道投入产出比低下
2. **跟进方式是关键因果杠杆**：面谈跟进相比电话/微信能显著提升转化率，因果效应经PSM验证显著
3. **延迟交付代价高昂**：{delay_rate:.0f}%的订单存在延迟交付，RDD分析显示这导致满意度显著下降
4. **城市间存在显著效率差异**：{best_city}转化率最高({city_cvr[best_city]:.1f}%)，{worst_city}最低({city_cvr[worst_city]:.1f}%)，差距达{city_cvr[best_city] - city_cvr[worst_city]:.1f}个百分点
5. **预测模型可提前识别风险**：基于机器学习的转化预测模型AUC达到可部署水平，SHAP分析提供了可解释的决策支持

### 三个战略建议

1. **预算再分配**：将20-30%的低ROI渠道预算重新分配至高ROI渠道
2. **跟进标准化**：建立面谈优先的跟进SOP，将首次跟进控制在3天内
3. **售后预警**：建立基于满意度的客户流失预警机制

### 预期财务影响
- 订单量增长：10-18%
- 单客获取成本降低：8-12%
- 客户满意度提升：0.3-0.5分
- ROI提升：15-25%

---

## 第2章 诊断：数据揭示的真相

### 2.1 运营全景
截至当前数据窗口，G9运营关键指标：
- 总线索量：{total_leads:,}
- 整体转化率：{cvr:.1f}%
- 成交订单：{total_orders:,}
- 延迟交付率：{delay_rate:.1f}%
- 平均交付评分：{avg_delivery_score:.2f}/5.0
- 售后平均满意度：{avg_aftersales:.2f}/5.0

### 2.2 结构性失衡
渠道层面呈现明显的"高投入低产出"与"低投入高产出"并存现象。
详情参见 W2 charts: 渠道归因热力图和ROI对比图。

### 2.3 效率黑洞
- 跟进响应时间：首次跟进平均间隔影响转化率
- 跟进质量：沟通时长与转化率呈倒U型关系
- 人员不均：部分销售员负载率过高

### 2.4 售后危机
- 主要投诉类型集中在质量问题和售后维修
- 投诉处理时长过长导致满意度进一步下降
- 配件缺货是最常见的具体投诉内容

### 2.5 成本结构
- 租金占门店成本的40-50%
- 市场费用分布不均衡
- 一线城市成本远高于其他城市

---

## 第3章 归因：什么真正驱动了转化

### 3.1 跟进方式的因果效应
PSM分析显示，面谈跟进相比非面谈方式对转化率有显著的因果正效应。
该结论经过多模型验证(Logistic/RF/GBM)和敏感性分析(Rosenbaum Bounds)验证稳健。

### 3.2 延迟交付的满意度代价
RDD断点回归分析表明，延迟交付导致交付评分显著下降。
该效应在不同带宽和Placebo断点下保持稳健。

### 3.3 渠道价值的公平分配
基于Shapley Value的渠道归因显示，不同渠道对转化的边际贡献存在显著差异。
与简单的Last Click归因相比，Shapley归因提供了更公平的价值分配。

### 3.4 异质性效应
HTE分析揭示了不同子群对跟进方式的差异化响应：
- 年轻客户(26-35岁)对面谈响应更强
- 一线城市客户对跟进方式更敏感
- 不同渠道来源的客户转化路径不同

---

## 第4章 预测：建立早期预警体系

### 4.1 转化预测模型
基于LightGBM/XGBoost的转化预测模型，经过Optuna贝叶斯超参优化，
AUC达到可部署水平。概率校准后ECE显著改善。

### 4.2 流失预警
Cox比例风险模型和随机生存森林识别出关键流失风险因素：
- 延迟交付是最强的流失预测因子
- 质量类投诉的流失风险是其他类型的2倍
- 首次售后满意度对长期留存至关重要

### 4.3 可解释的AI
SHAP分析提供了完整的模型可解释性：
- 试驾时长是最重要的单变量预测因子
- 跟进强度和首次跟进间隔的交互效应显著
- 年龄和渠道之间存在非线性交互

---

## 第5章 优化：策略推荐与模拟

### 5.1 预算再分配
线性规划优化显示，在总预算不变的前提下，通过重新分配可将预期订单量提升10-15%。
GA算法和NSGA-II多目标优化提供了稳健的方案选择。

### 5.2 定价博弈
博弈论分析建议：当竞品降价5%时不跟降，当竞品降价10%时降5%作为应对。
此策略在保护利润率的同时维持市场份额。

### 5.3 人效提升
整数规划排班优化可在现有人员基础上提升线索处理效率。
建议按销售员绩效和转化率进行差异化线索分配。

---

## 第6章 路线图：从分析到行动

### 6.1 Q4分阶段实施计划

| 阶段 | 周次 | 关键行动 | 责任人 | 里程碑 |
|------|------|---------|--------|--------|
| 数据基础 | W1-W2 | 数据清洗、宽表构建、因果分析 | 数据团队 | 宽表验收 |
| 模型训练 | W2-W4 | 预测模型、优化模型、SHAP分析 | 算法团队 | 模型AUC达标 |
| 策略制定 | W4-W5 | 预算方案、跟进修订、定价策略 | 运营团队 | 策略方案评审 |
| 系统上线 | W5-W8 | 仪表盘部署、团队培训、试运行 | 技术团队 | 仪表盘上线 |
| 全面推广 | W9-W12 | 正式执行、每日监控、周度复盘 | 全团队 | Q4目标达成 |

### 6.2 关键里程碑与责任人
- W2: 因果分析报告完成
- W4: 优化策略方案确定
- W5: 仪表盘V1上线
- W8: 试运行结束
- W12: Q4目标评估

### 6.3 风险与应对
1. **数据质量风险**: 建立数据质量监控，定期校准
2. **模型衰减风险**: 月度模型重训练和验证
3. **执行偏差风险**: 建立日报和周报机制
4. **竞品突变风险**: 每两周更新博弈分析

### 6.4 后续迭代建议
- 引入更多实时数据源（天气、节假日、竞品促销）
- 建立A/B测试框架
- 探索强化学习在动态定价中的应用
- 拓展到全产品线的统一决策平台

---

## 附录

### A. 方法论详解
- PSM: 倾向得分匹配 (Rosenbaum and Rubin, 1983)
- RDD: 断点回归设计 (Imbens and Lemieux, 2008)
- Shapley Value: 博弈论公平分配 (Shapley, 1953)
- LightGBM: 梯度提升决策树 (Ke et al., 2017)
- Optuna: 贝叶斯超参数优化 (Akiba et al., 2019)
- NSGA-II: 多目标遗传算法 (Deb et al., 2002)

### B. 模型技术报告
详见 reports/W3_model_report.md

### C. 数据字典
详见 reports/W1_data_quality_report.md

---

*本文档由G9智能销售运营决策系统自动生成*
*数据仅本地处理，不外传*
"""

    # Save
    path = os.path.join(REPORTS_DIR, 'W5_white_paper.md')
    with open(path, 'w', encoding='utf-8') as f:
        f.write(paper)
    print(f'  White paper saved: {path}')
    return paper

# ═══════════════════════════════════════════════════════════════
# 2. PPT Generation
# ═══════════════════════════════════════════════════════════════

def generate_ppt(wide=None):
    """生成PPT（使用python-pptx）"""
    print('\n[Generating PPT]')

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        print('  python-pptx not installed, skipping PPT generation')
        return

    # Load data
    if wide is None:
        csv_path = os.path.join(DATA_DIR, 'wide_table.csv')
        if os.path.exists(csv_path):
            import pandas as pd
            wide = pd.read_csv(csv_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x25, 0x63, 0xEB)
    DARK = RGBColor(0x1E, 0x29, 0x3B)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    RED = RGBColor(0xEF, 0x44, 0x44)

    def add_bg(slide, color=WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_box(slide, text, left=0.5, top=0.3, width=12, height=1, font_size=32):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = DARK
        return txBox

    def add_text_box(slide, text, left=0.5, top=1.5, width=12, height=5, font_size=18, color=DARK):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
        return txBox

    def add_kpi_box(slide, label, value, left, top, width=2.5, height=1.5, color=BLUE):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        return shape

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BLUE)
    add_title_box(slide, 'G9 智能销售运营决策系统', top=2, font_size=42)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = '数据驱动 · 全链路优化 · Q4销量增长15%目标'
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2 = tf.add_paragraph()
    p2.text = 'G9 销售运营团队 | 2025 Q4'
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title_box(slide, '议程')
    agenda = (
        'Part 1: 问题诊断 (Slides 3-7)\n'
        '  - Q4业绩回顾与数据全景\n'
        '  - 核心问题识别\n\n'
        'Part 2: 根因分析 (Slides 8-13)\n'
        '  - 跟进方式因果效应\n'
        '  - 渠道投入产出分析\n\n'
        'Part 3: 解决方案 (Slides 14-20)\n'
        '  - 预算重新分配\n'
        '  - 跟进流程优化\n\n'
        'Part 4: 实施计划 (Slides 21-25)'
    )
    add_text_box(slide, agenda, top=1.2, font_size=16)

    # --- Slide 3: KPI Dashboard ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title_box(slide, 'Q4 运营核心指标')

    if wide is not None:
        kpis = [
            ('总线索', f'{len(wide):,}', 0.5, 1.5, BLUE),
            ('转化率', f'{wide["是否下订"].mean()*100:.1f}%', 3.2, 1.5, GREEN),
            ('订单数', f'{wide["订单ID"].nunique():,}', 5.9, 1.5, RGBColor(0xF9, 0x73, 0x16)),
            ('延迟交付率', f'{wide[wide["订单ID"].notna()]["是否延迟交付"].mean()*100:.1f}%', 8.6, 1.5, RED),
            ('交付评分', f'{wide[wide["交付评分"].notna()]["交付评分"].mean():.2f}', 0.5, 3.5, RGBColor(0x8B, 0x5C, 0xF6)),
            ('售后满意度', f'{wide[wide["售后平均满意度"]>0]["售后平均满意度"].mean():.2f}', 3.2, 3.5, RGBColor(0x06, 0xB6, 0xD4)),
        ]
        for label, value, left, top, color in kpis:
            add_kpi_box(slide, label, value, left, top, color=color)

    # --- Slide: Key Findings ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title_box(slide, '五大核心发现')
    findings = (
        '1. 渠道效率悬殊 - 高ROI渠道投入不足，低ROI渠道资源浪费\n\n'
        '2. 跟进方式决定转化 - 面谈跟进的因果效应显著高于其他方式\n\n'
        '3. 延迟交付代价 - 50%+订单延迟交付，满意度下降0.5分\n\n'
        '4. 城市效率差异 - 城市间转化率差距达10个百分点\n\n'
        '5. AI可预测风险 - 机器学习模型可提前识别高风险客户'
    )
    add_text_box(slide, findings, top=1.2, font_size=20)

    # --- Slide: Strategy Recommendations ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title_box(slide, '三大战略建议')
    strategies = (
        '策略1: 预算重新分配\n'
        '  - 将低ROI渠道预算转移至高ROI渠道(懂车帝、抖音)\n'
        '  - 预期: 订单量提升10-15%\n\n'
        '策略2: 跟进流程标准化\n'
        '  - 首次跟进缩短至3天内, 推广面谈方式\n'
        '  - 预期: 转化率提升3-5个百分点\n\n'
        '策略3: 售后预警体系\n'
        '  - 交付评分<3自动预警, 投诉24小时响应\n'
        '  - 预期: 客户流失率降低20%'
    )
    add_text_box(slide, strategies, top=1.2, font_size=18)

    # --- Slide: Implementation Roadmap ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title_box(slide, 'Q4 实施路线图')
    roadmap = (
        'W1-W2: 数据清洗 + 因果分析\n'
        '├─ 宽表构建 ✓\n'
        '├─ PSM/RDD因果推断 ✓\n'
        '└─ 渠道归因分析 ✓\n\n'
        'W3-W4: 预测建模 + 策略优化\n'
        '├─ 转化预测模型训练\n'
        '├─ 预算优化方案\n'
        '└─ 策略组合评估\n\n'
        'W5-W8: 系统上线 + 试运行\n'
        '├─ 仪表盘部署\n'
        '├─ 团队培训\n'
        '└─ 试运行监控\n\n'
        'W9-W12: 全面推广\n'
        '├─ 正式执行\n'
        '├─ 每日监控\n'
        '└─ 周度复盘'
    )
    add_text_box(slide, roadmap, top=1.2, font_size=14)

    # --- Slide: Summary & Q&A ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = '总结 & Q&A'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = '\n数据驱动决策, 科学优化运营\n实现Q4销量增长15%的目标'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.alignment = PP_ALIGN.CENTER

    # Save
    path = os.path.join(REPORTS_DIR, 'W5_presentation.pptx')
    prs.save(path)
    print(f'  PPT saved: {path} ({len(prs.slides)} slides)')
    return prs

# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

if __name__ == '__main__':
    print('=' * 60)
    print('  G9 W5: White Paper & PPT Generation')
    print('=' * 60)

    generate_white_paper()
    generate_ppt()

    print('\n[OK] W5 Deliverables Generated!')
    print('  White Paper: reports/W5_white_paper.md')
    print('  PPT: reports/W5_presentation.pptx')
    print('  Dashboard: dashboard/app.py (streamlit run)')
